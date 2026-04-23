from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Apply a simple title slide
slide_layout_title = prs.slide_layouts[0]
slide_layout_bullet = prs.slide_layouts[1]

# Slide 1: Title
slide_1 = prs.slides.add_slide(slide_layout_title)
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "Serverless Object Detection Pipeline"
subtitle_1.text = "Production-ready AWS architecture using Amazon Rekognition\n\nProject Architecture & Implementation"

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(slide_layout_bullet)
title_2 = slide_2.shapes.title
title_2.text = "1. Introduction & Overview"
body_2 = slide_2.placeholders[1]
tf_2 = body_2.text_frame
tf_2.text = "What is this project?"
p = tf_2.add_paragraph()
p.text = "A powerful image analysis pipeline built on AWS Serverless."
p.level = 1
p = tf_2.add_paragraph()
p.text = "Leverages Amazon Rekognition to detect:"
p.level = 1
p = tf_2.add_paragraph()
p.text = "Objects and Labels (e.g., cars, buildings, animals)"
p.level = 2
p = tf_2.add_paragraph()
p.text = "Text in Images (OCR)"
p.level = 2
p = tf_2.add_paragraph()
p.text = "Moderation/Inappropriate Content"
p.level = 2
p = tf_2.add_paragraph()
p.text = "Personal Protective Equipment (PPE compliance)"
p.level = 2
p = tf_2.add_paragraph()
p.text = "Designed for high throughput and low latency end-to-end."
p.level = 1

# Slide 3: Architecture Overview
slide_3 = prs.slides.add_slide(slide_layout_bullet)
title_3 = slide_3.shapes.title
title_3.text = "2. Architecture Overview"
body_3 = slide_3.placeholders[1]
tf_3 = body_3.text_frame
tf_3.text = "End-to-End Serverless Flow:"
p = tf_3.add_paragraph()
p.text = "Client calls GET /presign via API Gateway -> Lambda."
p.level = 1
p = tf_3.add_paragraph()
p.text = "Client performs direct PUT upload of the image to AWS S3 using presigned URL."
p.level = 1
p = tf_3.add_paragraph()
p.text = "S3 ObjectCreated event triggers the main detect-objects Lambda."
p.level = 1
p = tf_3.add_paragraph()
p.text = "detect-objects runs parallel Amazon Rekognition queries."
p.level = 1
p = tf_3.add_paragraph()
p.text = "Results are stored in DynamoDB (DetectionResults table)."
p.level = 1
p = tf_3.add_paragraph()
p.text = "DynamoDB Streams triggers an annotation Lambda for bounding boxes."
p.level = 1

# Slide 4: Infrastructure as Code (IaC)
slide_4 = prs.slides.add_slide(slide_layout_bullet)
title_4 = slide_4.shapes.title
title_4.text = "3. Infrastructure & Deployment"
body_4 = slide_4.placeholders[1]
tf_4 = body_4.text_frame
tf_4.text = "Managed with AWS SAM (Serverless Application Model)"
p = tf_4.add_paragraph()
p.text = "Defined in template.yaml configuring all AWS resources."
p.level = 1
p = tf_4.add_paragraph()
p.text = "Key infrastructure components:"
p.level = 1
p = tf_4.add_paragraph()
p.text = "S3 bucket with lifecycle rules for original and annotated images."
p.level = 2
p = tf_4.add_paragraph()
p.text = "DynamoDB table with PAY_PER_REQUEST and 30-day TTL."
p.level = 2
p = tf_4.add_paragraph()
p.text = "SNS topic for Detection Alerts (high confidence notifications)."
p.level = 2
p = tf_4.add_paragraph()
p.text = "Deployment automated via deploy.sh and deploy.ps1 scripts."
p.level = 1

# Slide 5: Backend Implementation Details
slide_5 = prs.slides.add_slide(slide_layout_bullet)
title_5 = slide_5.shapes.title
title_5.text = "4. Backend & Parallel Processing"
body_5 = slide_5.placeholders[1]
tf_5 = body_5.text_frame
tf_5.text = "Optimized Python Lambda Functions"
p = tf_5.add_paragraph()
p.text = "detect-objects Lambda uses ThreadPoolExecutor max_workers=4."
p.level = 1
p = tf_5.add_paragraph()
p.text = "Parallel Rekognition API calls (Labels, Text, Moderation, PPE) preventing one round-trip from blocking others."
p.level = 2
p = tf_5.add_paragraph()
p.text = "Image Annotation Lambda leverages DynamoDB Streams insert events -> uses Pillow for async drawing over S3 images."
p.level = 1
p = tf_5.add_paragraph()
p.text = "Unit tests implemented via `pytest` with `boto3` fully mocked (no AWS credentials needed locally)."
p.level = 1

# Slide 6: Client Interfaces
slide_6 = prs.slides.add_slide(slide_layout_bullet)
title_6 = slide_6.shapes.title
title_6.text = "5. Multi-Client Ecosystem"
body_6 = slide_6.placeholders[1]
tf_6 = body_6.text_frame
tf_6.text = "The backend supports multiple client platforms via Presigned upload pattern:"
p = tf_6.add_paragraph()
p.text = "React Native App (rn-app/): Cross-platform mobile functionality connecting REST APIs."
p.level = 1
p = tf_6.add_paragraph()
p.text = "Native Android App (android/): Written natively in Kotlin for optimized performance and device integration natively."
p.level = 1
p = tf_6.add_paragraph()
p.text = "Web Frontend (frontend/): Standard web implementation utilizing HTML5 drag-and-drop and polling."
p.level = 1
p = tf_6.add_paragraph()
p.text = "All clients leverage the decoupled `GET /presign` and `GET /results` polling architecture."
p.level = 1

# Slide 7: Security & Best Practices
slide_7 = prs.slides.add_slide(slide_layout_bullet)
title_7 = slide_7.shapes.title
title_7.text = "6. Security & Performance Best Practices"
body_7 = slide_7.placeholders[1]
tf_7 = body_7.text_frame
tf_7.text = "Bypassing Lambda memory limits with Presigned URLs:"
p = tf_7.add_paragraph()
p.text = "Lambda never touches the binary avoiding API Gateway limits and minimizing execution costs."
p.level = 1
p = tf_7.add_paragraph()
p.text = "Strict CORS configurations and API restrictions."
p.level = 1
p = tf_7.add_paragraph()
p.text = "Least Privilege IAM Policies mapped explicitly using AWS SAM."
p.level = 1
p = tf_7.add_paragraph()
p.text = "Auto-cleanup of data:"
p.level = 1
p = tf_7.add_paragraph()
p.text = "DynamoDB TTL to automatically purge records past 30 days."
p.level = 2
p = tf_7.add_paragraph()
p.text = "S3 Lifecycle rule pruning."
p.level = 2

# Slide 8: Future Extensions & Conclusion
slide_8 = prs.slides.add_slide(slide_layout_bullet)
title_8 = slide_8.shapes.title
title_8.text = "7. Future Enhancements & Conclusion"
body_8 = slide_8.placeholders[1]
tf_8 = body_8.text_frame
tf_8.text = "Conclusion: A robust, cost-effective serverless computer vision solution."
p = tf_8.add_paragraph()
p.text = "Planned Improvements:"
p.level = 1
p = tf_8.add_paragraph()
p.text = "Dead-letter queues (SQS DLQ) for failed detections with retry policies."
p.level = 2
p = tf_8.add_paragraph()
p.text = "Dashboard insights utilizing CloudWatch dashboards (p99 latency, rate tracking)."
p.level = 2
p = tf_8.add_paragraph()
p.text = "WebSockets for real-time client push notifications, replacing HTTP polling."
p.level = 2
p = tf_8.add_paragraph()
p.text = "The repository presents a flexible foundation for scaling out ML detection use cases at almost zero idle cost."
p.level = 1

# Save the presentation
prs.save('Project_Presentation.pptx')
print('Presentation generated successfully as Project_Presentation.pptx')
